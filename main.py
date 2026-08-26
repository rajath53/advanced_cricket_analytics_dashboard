
from fastapi import FastAPI, UploadFile, File, HTTPException
from fastapi.middleware.cors import CORSMiddleware
from fastapi.responses import StreamingResponse
from pathlib import Path
import json, io, pandas as pd
from reportlab.lib.pagesizes import landscape, A4
from reportlab.pdfgen import canvas
from pptx import Presentation
from pptx.util import Inches, Pt

BASE = Path(__file__).resolve().parent
DATA = json.loads((BASE / "sample_data.json").read_text(encoding="utf-8"))

app = FastAPI(title="Cricket Analytics API", version="1.0.0")
app.add_middleware(
    CORSMiddleware,
    allow_origins=["*"],
    allow_credentials=True,
    allow_methods=["*"],
    allow_headers=["*"],
)

@app.get("/api/health")
def health():
    return {"status": "ok"}

@app.get("/api/match")
def match():
    return DATA

@app.get("/api/sessions")
def sessions(day: int | None = None):
    rows = DATA["sessions"]
    if day:
        rows = [x for x in rows if x["day"] == day]
    return rows

@app.get("/api/batting")
def batting(innings: int | None = None):
    rows = DATA["batting"]
    if innings:
        rows = [x for x in rows if x["innings"] == innings]
    return rows

@app.get("/api/bowling")
def bowling(innings: int | None = None):
    rows = DATA["bowling"]
    if innings:
        rows = [x for x in rows if x["innings"] == innings]
    return rows

@app.get("/api/players/{player}")
def player(player: str):
    bat = [x for x in DATA["batting"] if x["player"].lower() == player.lower()]
    bowl = [x for x in DATA["bowling"] if x["player"].lower() == player.lower()]
    return {"player": player, "batting": bat, "bowling": bowl}

@app.post("/api/upload/excel")
async def upload_excel(file: UploadFile = File(...)):
    if not file.filename.lower().endswith((".xlsx", ".xls", ".csv")):
        raise HTTPException(400, "Upload Excel or CSV only")
    raw = await file.read()
    try:
        if file.filename.lower().endswith(".csv"):
            df = pd.read_csv(io.BytesIO(raw))
        else:
            df = pd.read_excel(io.BytesIO(raw))
        return {
            "filename": file.filename,
            "rows": len(df),
            "columns": list(df.columns),
            "preview": df.head(20).fillna("").to_dict(orient="records"),
            "message": "File accepted. Add column-mapping rules for your preferred source format."
        }
    except Exception as e:
        raise HTTPException(400, f"Could not read file: {e}")

def build_pdf():
    buf = io.BytesIO()
    c = canvas.Canvas(buf, pagesize=landscape(A4))
    w, h = landscape(A4)
    c.setFont("Helvetica-Bold", 24)
    c.drawString(40, h-55, DATA["match"]["title"])
    c.setFont("Helvetica", 12)
    c.drawString(40, h-78, DATA["match"]["venue"] + " | " + DATA["match"]["dates"])
    y = h-120
    c.setFont("Helvetica-Bold", 15)
    c.drawString(40, y, "Match summary")
    y -= 25
    c.setFont("Helvetica", 11)
    for line in [
        DATA["match"]["result"],
        "Sri Lanka: 308 & 101",
        "West Indies: 626/9d",
        "Largest partnership: 401 — Jangoo + Chase",
        "Top score: Amir Jangoo 233",
        "Kemar Roach: 300th Test wicket"
    ]:
        c.drawString(55, y, line); y -= 18
    c.showPage()
    c.setFont("Helvetica-Bold", 18)
    c.drawString(40, h-50, "Session analysis")
    y = h-85
    c.setFont("Helvetica-Bold", 10)
    c.drawString(40, y, "Day")
    c.drawString(85, y, "Session")
    c.drawString(150, y, "Runs")
    c.drawString(205, y, "Wickets")
    c.drawString(275, y, "Overs")
    y -= 18
    c.setFont("Helvetica", 10)
    for s in DATA["sessions"]:
        c.drawString(40, y, str(s["day"]))
        c.drawString(85, y, str(s["session"]))
        c.drawString(150, y, str(s["runs"]))
        c.drawString(205, y, str(s["wickets"]))
        c.drawString(275, y, str(s["overs"]))
        y -= 16
        if y < 45:
            c.showPage(); y = h-50
    c.save()
    buf.seek(0)
    return buf

@app.get("/api/export/pdf")
def export_pdf():
    buf = build_pdf()
    return StreamingResponse(buf, media_type="application/pdf",
        headers={"Content-Disposition": "attachment; filename=cricket_match_report.pdf"})

@app.get("/api/export/pptx")
def export_pptx():
    prs = Presentation()
    slide = prs.slides.add_slide(prs.slide_layouts[0])
    slide.shapes.title.text = DATA["match"]["title"]
    slide.placeholders[1].text = DATA["match"]["result"]

    slide = prs.slides.add_slide(prs.slide_layouts[5])
    slide.shapes.title.text = "Match overview"
    tx = slide.shapes.add_textbox(Inches(0.7), Inches(1.4), Inches(11.5), Inches(4.8))
    tf = tx.text_frame
    for text in [
        "West Indies: 626/9d",
        "Sri Lanka: 308 & 101",
        "Largest partnership: 401",
        "Amir Jangoo: 233",
        "Roston Chase: 194",
        "Kemar Roach: 300th Test wicket"
    ]:
        p = tf.add_paragraph(); p.text = text; p.font.size = Pt(22)

    slide = prs.slides.add_slide(prs.slide_layouts[5])
    slide.shapes.title.text = "Session analysis"
    rows = len(DATA["sessions"]) + 1
    table = slide.shapes.add_table(rows, 5, Inches(0.6), Inches(1.3), Inches(11.8), Inches(5.2)).table
    headers = ["Day", "Session", "Runs", "Wickets", "Overs"]
    for i, h in enumerate(headers): table.cell(0,i).text = h
    for r, s in enumerate(DATA["sessions"], 1):
        vals = [s["day"], s["session"], s["runs"], s["wickets"], s["overs"]]
        for i, v in enumerate(vals): table.cell(r,i).text = str(v)
    out = io.BytesIO(); prs.save(out); out.seek(0)
    return StreamingResponse(out, media_type="application/vnd.openxmlformats-officedocument.presentationml.presentation",
        headers={"Content-Disposition": "attachment; filename=cricket_match_report.pptx"})
